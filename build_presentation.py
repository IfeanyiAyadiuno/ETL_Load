"""Generate ``presentation.pptx`` from ``presentation.md``.

Reads the slide outline (``## Slide N — Title`` sections separated by ``---``)
and renders a clean deck:

- Title slide with brand-green accent
- Content slides: title bar, bullet list on the left, screenshot placeholder
  box on the right, speaker notes from the ``**Speaker notes:**`` block

Usage:
    python build_presentation.py
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
SRC = ROOT / "presentation.md"
OUT = ROOT / "presentation.pptx"

BRAND_GREEN = RGBColor(0x1A, 0x4D, 0x3E)
ACCENT_GREEN = RGBColor(0x2E, 0x8B, 0x57)
LIGHT_GREY = RGBColor(0xF1, 0xF3, 0xF2)
DARK_TEXT = RGBColor(0x21, 0x2B, 0x2F)
PLACEHOLDER_BORDER = RGBColor(0xB7, 0xC7, 0xC2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


@dataclass
class Slide:
    title: str
    bullets: List[str] = field(default_factory=list)
    placeholder: Optional[str] = None
    notes: str = ""


_BOLD_RE = re.compile(r"\*\*([^*]+)\*\*")
_BACKTICK_RE = re.compile(r"`([^`]+)`")
_ITALIC_RE = re.compile(r"(?<!\w)_([^_]+)_(?!\w)")


def _strip_md(text: str) -> str:
    text = _BOLD_RE.sub(r"\1", text)
    text = _BACKTICK_RE.sub(r"\1", text)
    text = _ITALIC_RE.sub(r"\1", text)
    return text.strip()


def _extract_placeholder(line: str) -> str:
    text = re.sub(r"^\*\*[^*]+\*\*\s*:?\s*", "", line.strip())
    return _strip_md(text)


def parse_markdown(text: str) -> List[Slide]:
    blocks = [b.strip() for b in text.split("\n---\n")]
    slides: List[Slide] = []
    for block in blocks:
        m = re.search(r"^##\s+Slide\s+\d+\s+[—-]\s+(.+)$", block, re.MULTILINE)
        if not m:
            continue

        title = _strip_md(m.group(1))
        bullets: List[str] = []
        placeholder: Optional[str] = None
        notes_lines: List[str] = []
        in_notes = False
        in_placeholder = False

        for raw in block.splitlines():
            line = raw.rstrip()
            stripped = line.lstrip()

            if in_notes:
                notes_lines.append(line)
                continue

            if line.startswith("- "):
                bullets.append(_strip_md(line[2:].strip()))
                in_placeholder = False
                continue

            if re.match(r"^\s+[-*]\s", line) and bullets:
                bullets[-1] += " " + _strip_md(re.sub(r"^\s+[-*]\s+", "", line))
                continue

            if stripped.startswith("**Speaker notes:**"):
                in_notes = True
                in_placeholder = False
                continue

            if stripped.startswith("**Screenshot"):
                placeholder = _extract_placeholder(stripped)
                in_placeholder = True
                continue

            if in_placeholder and stripped and not stripped.startswith("**"):
                placeholder = (placeholder + " " + _strip_md(stripped)).strip()
                continue

            if not stripped:
                in_placeholder = False

        notes = "\n".join(notes_lines).strip()
        slides.append(Slide(title=title, bullets=bullets, placeholder=placeholder, notes=notes))

    return slides


def _set_run(run, text: str, *, size: int, bold: bool = False, color: RGBColor = DARK_TEXT) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs: Presentation, slide_data: Slide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(2.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_GREEN
    bar.line.fill.background()
    bar.shadow.inherit = False

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.6), sw, Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GREEN
    accent.line.fill.background()
    accent.shadow.inherit = False

    title_text = slide_data.bullets[0] if slide_data.bullets else slide_data.title

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.95), sw - Inches(1.4), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = None
    p.add_run()
    _set_run(p.runs[0], title_text, size=40, bold=True, color=BRAND_GREEN)

    for sub in slide_data.bullets[1:]:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.add_run()
        _set_run(p.runs[0], sub, size=20, color=DARK_TEXT)

    if slide_data.notes:
        slide.notes_slide.notes_text_frame.text = slide_data.notes


def add_content_slide(prs: Presentation, slide_data: Slide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_GREEN
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar_tf = bar.text_frame
    bar_tf.margin_left = Inches(0.45)
    bar_tf.margin_top = Inches(0.18)
    bar_tf.margin_right = Inches(0.4)
    bar_tf.margin_bottom = Inches(0.05)
    bp = bar_tf.paragraphs[0]
    bp.add_run()
    _set_run(bp.runs[0], slide_data.title, size=26, bold=True, color=WHITE)

    body_top = Inches(1.25)
    body_height = sh - body_top - Inches(0.45)

    has_placeholder = bool(slide_data.placeholder)
    if has_placeholder:
        bullet_left = Inches(0.55)
        bullet_width = Inches(7.0)
        ph_left = Inches(7.85)
        ph_width = sw - ph_left - Inches(0.45)
    else:
        bullet_left = Inches(0.7)
        bullet_width = sw - Inches(1.4)

    bb = slide.shapes.add_textbox(bullet_left, body_top, bullet_width, body_height)
    bf = bb.text_frame
    bf.word_wrap = True

    for i, bullet in enumerate(slide_data.bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.space_after = Pt(8)
        p.add_run()
        _set_run(p.runs[0], "•  " + bullet, size=18, color=DARK_TEXT)

    if has_placeholder:
        ph = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, ph_left, body_top, ph_width, body_height
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_GREY
        ph.line.color.rgb = PLACEHOLDER_BORDER
        ph.line.width = Pt(1.25)
        ph.shadow.inherit = False
        pf = ph.text_frame
        pf.word_wrap = True
        pf.margin_left = Inches(0.3)
        pf.margin_right = Inches(0.3)
        pf.margin_top = Inches(0.4)
        pp = pf.paragraphs[0]
        pp.add_run()
        _set_run(pp.runs[0], "[ SCREENSHOT ]", size=15, bold=True, color=ACCENT_GREEN)
        cap = pf.add_paragraph()
        cap.space_before = Pt(10)
        cap.add_run()
        _set_run(cap.runs[0], slide_data.placeholder, size=13, color=DARK_TEXT)

    if slide_data.notes:
        slide.notes_slide.notes_text_frame.text = slide_data.notes


def main() -> None:
    text = SRC.read_text(encoding="utf-8")
    slides = parse_markdown(text)
    if not slides:
        raise SystemExit(f"No slides found in {SRC}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        if i == 0:
            add_title_slide(prs, slide_data)
        else:
            add_content_slide(prs, slide_data)

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(slides)} slides)")


if __name__ == "__main__":
    main()
