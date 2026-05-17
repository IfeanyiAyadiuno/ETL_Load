"""
Build ``docs/presentation.pptx`` from this folder.

Run::

    python -m pip install python-pptx
    python docs/build_presentation.py

The deck is intentionally generated from code (rather than hand-built in
PowerPoint) so the slide deck stays in lock-step with project messaging
and can be regenerated whenever the talking points change.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Theme — Pacific Canbriam: blue / green / white
# ---------------------------------------------------------------------------

BLUE = RGBColor(0x1A, 0x2C, 0x5C)          # navy from PCE logo wordmark
BLUE_DEEP = RGBColor(0x10, 0x1F, 0x42)     # darkest navy (title-slide hero)
BLUE_BRIGHT = RGBColor(0x35, 0x77, 0xB5)   # mid-blue from PCE logo orb
GREEN = RGBColor(0x82, 0xBC, 0x3F)         # leaf green from PCE logo
GREEN_DEEP = RGBColor(0x6A, 0x9E, 0x2E)    # darker leaf green
SURFACE_BLUE = RGBColor(0xE8, 0xEF, 0xF8)  # light blue-tint surface
SURFACE_GREEN = RGBColor(0xEF, 0xF6, 0xE2) # light green-tint surface
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x29, 0x37)           # body text
MUTED = RGBColor(0x6B, 0x72, 0x80)
DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Calibri"
FONT_DISPLAY = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ORG_NAME = "Pacific Canbriam Energy"
PROJECT_NAME = "Production Update System"


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def set_solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def set_line(shape, rgb: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    rect = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is not None:
        set_solid_fill(rect, fill)
    else:
        rect.fill.background()
    if line is not None:
        set_line(rect, line, 0.75)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font=FONT,
    size=18,
    bold=False,
    italic=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_paragraphs(
    slide,
    x,
    y,
    w,
    h,
    paragraphs,
    *,
    font=FONT,
    size=16,
    color=INK,
    align=PP_ALIGN.LEFT,
    line_spacing=1.25,
    space_after=6,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True

    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)

        for j, segment in enumerate(item if isinstance(item, list) else [item]):
            if isinstance(segment, str):
                segment = {"text": segment}
            run = p.add_run()
            run.text = segment.get("text", "")
            run.font.name = segment.get("font", font)
            run.font.size = Pt(segment.get("size", size))
            run.font.bold = segment.get("bold", False)
            run.font.italic = segment.get("italic", False)
            run.font.color.rgb = segment.get("color", color)
    return tb


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items,
    *,
    bullet_color=GREEN,
    text_color=INK,
    size=16,
    bold_lead=True,
):
    """
    Render bullet points where each item is either:
      - "Plain bullet text"
      - ("Bold lead", "continuation text")
    Bullet marker is a small gold square; no Office bullet formatting.
    """
    bullet_size = Inches(0.10)
    line_gap = Inches(0.62) if size >= 18 else Inches(0.50)
    text_x = x + bullet_size + Inches(0.18)
    text_w = w - (bullet_size + Inches(0.18))

    cursor = y
    for item in items:
        if isinstance(item, tuple):
            lead, rest = item
        else:
            lead, rest = item, ""

        bullet = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x,
            cursor + Inches(0.18),
            bullet_size,
            bullet_size,
        )
        set_solid_fill(bullet, bullet_color)

        segments = []
        if bold_lead and rest:
            segments.append({"text": lead, "bold": True, "color": text_color, "size": size})
            segments.append({"text": " — " + rest, "color": text_color, "size": size})
        elif rest:
            segments.append({"text": lead + " — ", "bold": True, "color": text_color, "size": size})
            segments.append({"text": rest, "color": text_color, "size": size})
        else:
            segments.append({"text": lead, "color": text_color, "size": size,
                             "bold": bold_lead})

        add_paragraphs(slide, text_x, cursor, text_w, line_gap, [segments],
                       size=size, color=text_color, line_spacing=1.2,
                       space_after=0)
        cursor += line_gap
    return cursor


# ---------------------------------------------------------------------------
# Slide chrome (title slide, content slide, section divider)
# ---------------------------------------------------------------------------


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_full_background(slide, rgb):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=rgb)


def title_slide(prs, title, subtitle, footer_left, footer_right):
    slide = slide_blank(prs)
    add_full_background(slide, BLUE_DEEP)

    add_rect(slide, 0, 0, Inches(0.45), SLIDE_H, fill=GREEN)

    add_text(slide, Inches(0.95), Inches(2.6), Inches(11.5), Inches(0.5),
             ORG_NAME.upper(), size=14, bold=True, color=GREEN,
             align=PP_ALIGN.LEFT)

    add_text(slide, Inches(0.95), Inches(3.05), Inches(11.5), Inches(1.4),
             title, font=FONT_DISPLAY, size=54, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT)

    add_rect(slide, Inches(0.95), Inches(4.45), Inches(0.7), Pt(2), fill=GREEN)

    add_text(slide, Inches(0.95), Inches(4.65), Inches(11.5), Inches(0.8),
             subtitle, size=20, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, Inches(0.95), Inches(6.85), Inches(6),  Inches(0.4),
             footer_left, size=12, color=GREEN, bold=True)
    add_text(slide, Inches(7.0), Inches(6.85), Inches(5.4), Inches(0.4),
             footer_right, size=12, color=WHITE, align=PP_ALIGN.RIGHT)
    return slide


def content_slide(prs, title, eyebrow=None, slide_no=None, total=None):
    slide = slide_blank(prs)
    add_full_background(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=BLUE)
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Pt(3), fill=GREEN)

    if eyebrow:
        add_text(slide, Inches(0.6), Inches(0.18), Inches(8), Inches(0.35),
                 eyebrow.upper(), size=11, bold=True, color=GREEN)
        title_y = Inches(0.50)
    else:
        title_y = Inches(0.30)

    add_text(slide, Inches(0.6), title_y, Inches(11), Inches(0.85),
             title, font=FONT_DISPLAY, size=30, bold=True, color=WHITE)

    if slide_no is not None and total is not None:
        add_text(slide, Inches(11.6), Inches(0.30), Inches(1.4), Inches(0.75),
                 f"{slide_no:02d} / {total:02d}", size=14, bold=True,
                 color=GREEN, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(slide, Inches(0.6), Inches(7.10), SLIDE_W - Inches(1.2), Pt(0.75),
             fill=DIVIDER)
    add_text(slide, Inches(0.6), Inches(7.18), Inches(8), Inches(0.25),
             f"{PROJECT_NAME}  ·  {ORG_NAME}", size=10, color=MUTED)
    return slide


def section_divider(prs, number, kicker, title, slide_no=None, total=None):
    slide = slide_blank(prs)
    add_full_background(slide, BLUE)

    add_rect(slide, 0, Inches(3.10), SLIDE_W, Pt(2), fill=GREEN)

    add_text(slide, Inches(0.8), Inches(2.30), Inches(4), Inches(0.85),
             f"{number:02d}", font=FONT_DISPLAY, size=80, bold=True,
             color=GREEN, align=PP_ALIGN.LEFT)

    add_text(slide, Inches(0.8), Inches(3.35), Inches(11), Inches(0.45),
             kicker.upper(), size=13, bold=True, color=GREEN)

    add_text(slide, Inches(0.8), Inches(3.80), Inches(11.7), Inches(1.3),
             title, font=FONT_DISPLAY, size=40, bold=True, color=WHITE)

    if slide_no is not None and total is not None:
        add_text(slide, Inches(11.6), Inches(6.85), Inches(1.4), Inches(0.4),
                 f"{slide_no:02d} / {total:02d}", size=11, bold=True,
                 color=GREEN, align=PP_ALIGN.RIGHT)
    return slide


# ---------------------------------------------------------------------------
# Composite components
# ---------------------------------------------------------------------------


def add_card(slide, x, y, w, h, *, title, items, accent=GREEN, fill=SURFACE_BLUE,
             title_color=BLUE, body_color=INK, title_size=18, body_size=14):
    add_rect(slide, x, y, w, h, fill=fill, line=DIVIDER)
    add_rect(slide, x, y, Inches(0.18), h, fill=accent)

    add_text(slide, x + Inches(0.55), y + Inches(0.30),
             w - Inches(0.75), Inches(0.50),
             title, font=FONT_DISPLAY, size=title_size, bold=True,
             color=title_color)

    add_bullets(slide, x + Inches(0.55), y + Inches(0.95),
                w - Inches(0.85), h - Inches(1.1), items,
                bullet_color=accent, text_color=body_color, size=body_size,
                bold_lead=False)


def add_stat(slide, x, y, w, h, *, value, label, accent=GREEN):
    add_rect(slide, x, y, w, h, fill=BLUE)
    add_rect(slide, x, y, w, Pt(3), fill=accent)
    add_text(slide, x, y + Inches(0.35), w, Inches(1.2),
             value, font=FONT_DISPLAY, size=44, bold=True, color=accent,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + h - Inches(0.85), w, Inches(0.6),
             label.upper(), size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_arrow(slide, x, y, w, h, fill=GREEN):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    set_solid_fill(arrow, fill)
    return arrow


def add_table_card(slide, x, y, w, h, headers, rows,
                   header_fill=BLUE, header_text=WHITE,
                   alt_fill=SURFACE_GREEN, body_text=INK,
                   header_size=14, body_size=13):
    """Hand-rendered table-style card — one rectangle per cell so we can
    fully control colors and fonts (python-pptx tables are limited)."""
    n_cols = len(headers)
    col_widths_emu = [int(w / n_cols) for _ in range(n_cols)]
    diff = w - sum(col_widths_emu)
    col_widths_emu[-1] += diff

    n_rows = len(rows) + 1
    header_h = Inches(0.55)
    body_h_total = h - header_h
    row_h = int(body_h_total / len(rows))

    cx = x
    for i, header in enumerate(headers):
        cw = col_widths_emu[i]
        add_rect(slide, cx, y, cw, header_h, fill=header_fill)
        add_text(slide, cx + Inches(0.18), y, cw - Inches(0.36), header_h,
                 header, size=header_size, bold=True, color=header_text,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += cw

    for r, row in enumerate(rows):
        ry = y + header_h + r * row_h
        cx = x
        fill = WHITE if r % 2 == 0 else alt_fill
        for i, cell in enumerate(row):
            cw = col_widths_emu[i]
            add_rect(slide, cx, ry, cw, row_h, fill=fill)
            add_text(slide, cx + Inches(0.18), ry, cw - Inches(0.36), row_h,
                     cell, size=body_size, bold=(i == 0), color=body_text,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += cw


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def build_title(prs):
    title_slide(
        prs,
        title="Production Update System",
        subtitle="A centralized data platform for daily and monthly production.",
        footer_left=ORG_NAME.upper(),
        footer_right="Engineering Overview · 2026",
    )


def build_problem(prs, n, total):
    slide = content_slide(prs, "Everything lived in Excel.",
                          eyebrow="The problem we were solving",
                          slide_no=n, total=total)

    add_text(slide, Inches(0.6), Inches(1.65), Inches(12), Inches(0.6),
             "No central database. Every task was its own workbook.",
             size=20, italic=True, color=BLUE_BRIGHT)

    items = [
        "Manual copy-paste between Snowflake, ValNav, Accumap, and spreadsheets",
        "Allocation factors and sales ratios derived by hand — easy to break",
        "Re-running a month was risky; results varied between engineers",
        "Onboarding meant \u201Clearn the spreadsheets,\u201D not \u201Clearn the data\u201D",
        "No reliable way to audit how a number was produced",
    ]
    add_bullets(slide, Inches(0.6), Inches(2.55), Inches(7.4), Inches(4.0),
                items, size=17, bold_lead=False)

    add_card(slide, Inches(8.4), Inches(2.40), Inches(4.4), Inches(4.2),
             title="Common failure modes",
             items=[
                 "Two engineers editing copies of the same workbook",
                 "Mis-pasted columns silently changing totals",
                 "Lost public-sales data when re-doing PA",
                 "Different answers for the same month",
             ],
             accent=GREEN, fill=SURFACE_BLUE, title_size=16, body_size=13)


def build_what_we_built(prs, n, total):
    slide = content_slide(prs, "Two deliverables, designed together.",
                          eyebrow="What we built",
                          slide_no=n, total=total)

    add_text(slide, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
             "A new SQL Server data model — and the application that keeps it current.",
             size=18, color=MUTED)

    add_card(slide, Inches(0.6), Inches(2.40), Inches(5.8), Inches(4.4),
             title="SQL Server data model",
             items=[
                 "PCE_WM — Well Master",
                 "PCE_CDA — daily allocation",
                 "PCE_Production — daily production with sequences",
                 "Allocation_Factors — PA & Public Sales factors",
                 "PCE_Surveys, PCE_TC — surveys and type curves",
             ],
             accent=BLUE_BRIGHT, fill=SURFACE_BLUE,
             title_size=20, body_size=14)

    add_arrow(slide, Inches(6.55), Inches(4.20), Inches(0.4), Inches(0.6))

    add_card(slide, Inches(7.05), Inches(2.40), Inches(5.8), Inches(4.4),
             title="PyQt5 desktop application",
             items=[
                 "One window, one operation per task",
                 "Loads from Snowflake, ValNav, Accumap, Excel",
                 "Writes back to the SQL tables — logged & repeatable",
                 "Background workers; cancel & summary on every run",
             ],
             accent=GREEN, fill=SURFACE_GREEN,
             title_size=20, body_size=14)


def build_data_flow(prs, n, total):
    slide = content_slide(prs, "How data flows.",
                          eyebrow="Architecture at a glance",
                          slide_no=n, total=total)

    sources = [
        ("Snowflake / Prodview", "Daily wellhead, gathered, pressures"),
        ("ValNav", "Monthly S2 / condensate factors"),
        ("Accumap", "Monthly public sales gas"),
        ("Excel imports", "Surveys & type curves"),
    ]

    box_w = Inches(2.8)
    box_h = Inches(1.05)
    gap = Inches(0.20)
    start_x = Inches(0.6)
    base_y = Inches(2.0)
    for i, (name, sub) in enumerate(sources):
        y = base_y + i * (box_h + gap)
        add_rect(slide, start_x, y, box_w, box_h, fill=SURFACE_BLUE, line=DIVIDER)
        add_rect(slide, start_x, y, Inches(0.12), box_h, fill=BLUE_BRIGHT)
        add_text(slide, start_x + Inches(0.30), y + Inches(0.18),
                 box_w - Inches(0.40), Inches(0.40),
                 name, size=15, bold=True, color=BLUE)
        add_text(slide, start_x + Inches(0.30), y + Inches(0.55),
                 box_w - Inches(0.40), Inches(0.40),
                 sub, size=12, color=MUTED)

    arrow_x = start_x + box_w + Inches(0.10)
    add_arrow(slide, arrow_x, Inches(3.65), Inches(0.9), Inches(0.85))

    app_x = arrow_x + Inches(1.05)
    app_w = Inches(2.4)
    add_rect(slide, app_x, Inches(3.10), app_w, Inches(2.0), fill=BLUE)
    add_rect(slide, app_x, Inches(3.10), app_w, Pt(3), fill=GREEN)
    add_text(slide, app_x, Inches(3.30), app_w, Inches(0.55),
             "PRODUCTION UPDATE", size=12, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)
    add_text(slide, app_x, Inches(3.85), app_w, Inches(0.55),
             "APPLICATION", size=12, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)
    add_text(slide, app_x, Inches(4.45), app_w, Inches(0.55),
             "PyQt5 · ETL · audit log", size=12, color=WHITE,
             align=PP_ALIGN.CENTER)

    arrow2_x = app_x + app_w + Inches(0.10)
    add_arrow(slide, arrow2_x, Inches(3.65), Inches(0.9), Inches(0.85))

    sql_x = arrow2_x + Inches(1.05)
    sql_w = SLIDE_W - sql_x - Inches(0.6)
    add_card(slide, sql_x, Inches(2.0), sql_w, Inches(4.7),
             title="SQL Server",
             items=[
                 "PCE_WM",
                 "PCE_CDA",
                 "PCE_Production",
                 "Allocation_Factors",
                 "PCE_Surveys",
                 "PCE_TC",
             ],
             accent=GREEN, fill=SURFACE_BLUE, title_size=18, body_size=14)


def build_eight_operations(prs, n, total):
    slide = content_slide(prs, "Eight operations, one window.",
                          eyebrow="What the application does",
                          slide_no=n, total=total)

    ops = [
        ("Well Master", "Maintain PCE_WM; import new wells from Snowflake"),
        ("Prodview / Snowflake", "Daily production retrieve into CDA & Production"),
        ("PA Allocations", "Monthly ValNav factors → Allocation_Factors"),
        ("Public Sales", "Accumap sales gas + sales ratios"),
        ("Survey Import", "Directional / log surveys into PCE_Surveys"),
        ("Type Curves Import", "Excel type curves → PCE_TC → Production"),
        ("Whitson+ Mass Upload", "Placeholder for next phase"),
        ("Exports / Reports", "Placeholder for next phase"),
    ]

    cols = 4
    rows = 2
    grid_x = Inches(0.6)
    grid_y = Inches(1.85)
    grid_w = SLIDE_W - Inches(1.2)
    grid_h = Inches(5.0)
    gap = Inches(0.20)
    cell_w = (grid_w - gap * (cols - 1)) / cols
    cell_h = (grid_h - gap * (rows - 1)) / rows

    for i, (name, sub) in enumerate(ops):
        r, c = divmod(i, cols)
        x = grid_x + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        is_placeholder = name in ("Whitson+ Mass Upload", "Exports / Reports")
        accent = MUTED if is_placeholder else GREEN
        title_color = MUTED if is_placeholder else BLUE
        body_color = MUTED if is_placeholder else INK
        bg = SURFACE_BLUE if not is_placeholder else RGBColor(0xF1, 0xF2, 0xF4)

        add_rect(slide, x, y, cell_w, cell_h, fill=bg, line=DIVIDER)
        add_rect(slide, x, y, cell_w, Pt(3), fill=accent)
        add_text(slide, x + Inches(0.30), y + Inches(0.30),
                 cell_w - Inches(0.6), Inches(0.5),
                 name, size=16, bold=True, color=title_color)
        add_text(slide, x + Inches(0.30), y + Inches(0.85),
                 cell_w - Inches(0.6), Inches(1.4),
                 sub, size=12, color=body_color)


def build_daily(prs, n, total):
    slide = content_slide(prs, "Daily flow: Prodview / Snowflake.",
                          eyebrow="Snowflake → CDA + Production rebuild",
                          slide_no=n, total=total)

    steps = [
        ("Pick window", "Auto: rolling 18 months, ends today − 2 days"),
        ("Pull Snowflake", "ECF, Gas WH, CGR / water, WGR, pressures, allocation"),
        ("Refresh CDA", "Replace matching PCE_CDA rows for the window"),
        ("Rebuild Production", "Sequences, cumulatives, monthly averages"),
        ("Sync type curves", "PCE_TC materialized into PCE_Production"),
    ]

    step_w = Inches(2.30)
    step_h = Inches(2.6)
    gap = Inches(0.20)
    n_steps = len(steps)
    total_w = step_w * n_steps + gap * (n_steps - 1)
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.40)

    for i, (title, sub) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        add_rect(slide, x, y, step_w, step_h, fill=SURFACE_BLUE, line=DIVIDER)
        add_rect(slide, x, y, step_w, Inches(0.55), fill=BLUE)
        add_text(slide, x + Inches(0.20), y + Inches(0.06),
                 step_w - Inches(0.4), Inches(0.45),
                 f"STEP {i + 1}", size=11, bold=True, color=GREEN)
        add_text(slide, x + Inches(0.20), y + Inches(0.75),
                 step_w - Inches(0.4), Inches(0.75),
                 title, size=16, bold=True, color=BLUE)
        add_text(slide, x + Inches(0.20), y + Inches(1.45),
                 step_w - Inches(0.4), Inches(1.0),
                 sub, size=12, color=INK)

    add_text(slide, Inches(0.6), Inches(5.55), SLIDE_W - Inches(1.2),
             Inches(0.45),
             "One click. One log. One summary line with row counts and elapsed time.",
             size=16, italic=True, color=BLUE_BRIGHT, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.6), Inches(6.10), SLIDE_W - Inches(1.2),
             Inches(0.45),
             "Exception wells skipped automatically · type-curve & YE2 rows protected",
             size=13, color=MUTED, align=PP_ALIGN.CENTER)


def build_monthly(prs, n, total):
    slide = content_slide(prs, "Monthly flow: PA + Public Sales.",
                          eyebrow="ValNav and Accumap, the right way",
                          slide_no=n, total=total)

    items = [
        ("PA Allocations (ValNav)", "Writes monthly factors to Allocation_Factors"),
        ("Public Sales (Accumap)", "Merges sales gas; recomputes sales CGR"),
        ("Auto-propagation", "Allocation changes flow to PCE_CDA and PCE_Production"),
        ("Sales_Gas preserved", "Re-running PA does not erase Public Sales data"),
        ("Idempotent month", "Existing month is deleted and rebuilt consistently"),
    ]
    add_bullets(slide, Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.5),
                items, size=16)

    add_card(slide, Inches(8.4), Inches(2.0), Inches(4.4), Inches(4.6),
             title="The old foot-gun",
             items=[
                 "Doing PA after Public Sales used to wipe Sales_Gas to zero",
                 "Engineers had to restore Accumap data by hand",
                 "Now: PA reads existing Sales_Gas first, then writes",
                 "Public Sales data survives PA reloads",
             ],
             accent=GREEN, fill=SURFACE_GREEN, title_size=16, body_size=13,
             title_color=BLUE)


def build_supporting(prs, n, total):
    slide = content_slide(prs, "Well Master, Surveys, Type Curves.",
                          eyebrow="Supporting flows — now routine",
                          slide_no=n, total=total)

    add_card(slide, Inches(0.6), Inches(2.0), Inches(4.0), Inches(4.6),
             title="Well Master",
             items=[
                 "Snowflake-driven new well preview",
                 "Safe edits with composite-name handling",
                 "Honored everywhere — exception flag respected",
             ],
             accent=GREEN, fill=SURFACE_BLUE, title_size=18, body_size=13)

    add_card(slide, Inches(4.85), Inches(2.0), Inches(4.0), Inches(4.6),
             title="Surveys",
             items=[
                 "Bulk and mapped Excel layouts",
                 "Append vs overwrite behavior is explicit",
                 "Unmatched UWIs reported as a CSV",
             ],
             accent=BLUE_BRIGHT, fill=SURFACE_BLUE, title_size=18, body_size=13)

    add_card(slide, Inches(9.10), Inches(2.0), Inches(3.7), Inches(4.6),
             title="Type Curves",
             items=[
                 "Excel → PCE_TC",
                 "Synced into PCE_Production at ImportDate",
                 "Protected on every rebuild",
             ],
             accent=GREEN, fill=SURFACE_GREEN, title_size=18, body_size=13)


def build_time_savings(prs, n, total):
    slide = content_slide(prs, "How this saves engineers time.",
                          eyebrow="Before vs after",
                          slide_no=n, total=total)

    headers = ["Task", "Before — Excel only", "After — App + SQL"]
    rows = [
        ("Daily Snowflake → CDA → Production",
         "Half-day spreadsheet round-trip",
         "One dialog, hands-off, logged summary"),
        ("Monthly PA (ValNav)",
         "Manual factor calc + manual propagation",
         "One dialog; CDA / Production update automatically"),
        ("Public Sales (Accumap)",
         "Manual UWI matching; risk of erasing PA",
         "One dialog; unmatched UWIs reported; PA preserved"),
        ("New well onboarding",
         "Update multiple workbooks",
         "Add in Well Master, run Prodview when ready"),
        ("Survey / Type Curve imports",
         "Manual mapping and pasting",
         "Bulk + mapped imports with audit reports"),
        ("Re-running a month",
         "Risky; could overwrite",
         "Safe and idempotent"),
    ]

    add_table_card(slide, Inches(0.6), Inches(1.85),
                   SLIDE_W - Inches(1.2), Inches(4.5),
                   headers, rows, header_size=14, body_size=12)

    add_text(slide, Inches(0.6), Inches(6.55),
             SLIDE_W - Inches(1.2), Inches(0.45),
             "Engineers stop doing data plumbing — and spend time on interpretation.",
             size=16, italic=True, bold=True, color=BLUE_BRIGHT,
             align=PP_ALIGN.CENTER)


def build_reliability(prs, n, total):
    slide = content_slide(prs, "Reliability and auditability.",
                          eyebrow="Built for engineering trust",
                          slide_no=n, total=total)

    items_left = [
        "Header, step lines, and summary on every run",
        "Row counts and elapsed time recorded",
        "Cancel supported during long jobs",
        "Date caps prevent writing into incomplete days",
    ]
    items_right = [
        "Exception-flagged wells skipped automatically",
        "Type-curve and YE2 rows protected during rebuilds",
        "Concurrent engineers share the same SQL tables",
        "No silent overwrites; nothing happens off-screen",
    ]

    add_bullets(slide, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5),
                items_left, size=16)
    add_bullets(slide, Inches(7.0), Inches(2.0), Inches(6.0), Inches(4.5),
                items_right, size=16)

    add_rect(slide, Inches(0.6), Inches(6.30), SLIDE_W - Inches(1.2),
             Inches(0.55), fill=SURFACE_BLUE)
    add_text(slide, Inches(0.6), Inches(6.30), SLIDE_W - Inches(1.2),
             Inches(0.55),
             "If something goes wrong, you can read the log and reproduce the result.",
             size=14, italic=True, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_next(prs, n, total):
    slide = content_slide(prs, "What's next.",
                          eyebrow="Roadmap — shipped value, clear next phase",
                          slide_no=n, total=total)

    items = [
        ("Whitson+ Mass Upload", "Replace the current placeholder workflow"),
        ("Exports / Reports", "First-class dialog and scheduled extracts"),
        ("Automated reconciliation", "Cross-source checks (Snowflake / Accumap / ValNav)"),
        ("Smaller asks", "Better progress signals, more name-resolution rules"),
    ]
    add_bullets(slide, Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.4),
                items, size=17)

    add_card(slide, Inches(8.9), Inches(2.0), Inches(4.0), Inches(4.4),
             title="Owners & timelines",
             items=[
                 "Tech lead — to be assigned",
                 "Reservoir lead — to be assigned",
                 "Pilot users — volunteers welcome",
             ],
             accent=BLUE_BRIGHT, fill=SURFACE_GREEN, title_size=16, body_size=13)


def build_qa(prs, n, total):
    slide = content_slide(prs, "Questions.",
                          eyebrow="Discussion",
                          slide_no=n, total=total)

    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(1.0),
             "Thanks for your time.", font=FONT_DISPLAY, size=44, bold=True,
             color=BLUE)

    add_rect(slide, Inches(0.6), Inches(3.45), Inches(0.7), Pt(2), fill=GREEN)

    add_text(slide, Inches(0.6), Inches(3.65), Inches(12), Inches(0.5),
             "Two minutes of Q&A — or grab me after.", size=20, color=BLUE_BRIGHT)

    add_card(slide, Inches(0.6), Inches(4.55), Inches(6.0), Inches(2.1),
             title="Where to file requests / bugs",
             items=[
                 "<email or ticket link>",
                 "Tag: Production Update System",
             ],
             accent=GREEN, fill=SURFACE_BLUE, title_size=16, body_size=13)

    add_card(slide, Inches(6.85), Inches(4.55), Inches(6.0), Inches(2.1),
             title="Documentation",
             items=[
                 "docs/USER_GUIDE.md — operators",
                 "docs/DEV_GUIDE.md — overview",
                 "docs/DEV_GUIDE_LAYOUT.md — deep technical map",
             ],
             accent=BLUE_BRIGHT, fill=SURFACE_BLUE, title_size=16, body_size=13)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        build_problem,
        build_what_we_built,
        build_data_flow,
        build_eight_operations,
        build_daily,
        build_monthly,
        build_supporting,
        build_time_savings,
        build_reliability,
        build_next,
        build_qa,
    ]
    total = len(builders) + 1

    build_title(prs)
    for idx, builder in enumerate(builders, start=2):
        builder(prs, idx, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path} ({total} slides)")


if __name__ == "__main__":
    here = Path(__file__).resolve().parent
    build(here / "presentation.pptx")
